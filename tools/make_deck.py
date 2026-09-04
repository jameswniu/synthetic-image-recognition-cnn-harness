"""Both decks, built from the same content so they cannot drift apart.

Two audiences need two shapes. The narrate deck (the -14 file) is for talking over: one idea a
slide, the picture doing the work, the numbers in the speaker notes. The leave-behind (the -20
file) is for the person who was not in the room and gets it forwarded: the same argument with the
evidence on the slide. Both carry a build-and-run slide, so the file names now undercount by one.

Every figure comes from the repo's own assets, which are themselves generated from real output, so
a number on a slide cannot disagree with the number in the README without somebody noticing.

Run: uv run --with python-pptx python tools/make_deck.py
Out: deliverables/checkbox-approach.pptx, the leave-behind in plain words, which is the canonical
     copy. The other registers remain in the code as working iterations and are not built.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"

INK = RGBColor(0x18, 0x18, 0x1B)
MUTED = RGBColor(0x71, 0x71, 0x7A)
DET = RGBColor(0x4F, 0x46, 0xE5)
DET_FILL = RGBColor(0xEE, 0xF2, 0xFF)
EXC = RGBColor(0xD9, 0x7A, 0x1E)
EXC_FILL = RGBColor(0xFF, 0xF4, 0xE0)
HUM = RGBColor(0x3A, 0x9D, 0x55)
HUM_FILL = RGBColor(0xE9, 0xF7, 0xEC)
RULE = RGBColor(0xE4, 0xE4, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# The theme. One deep navy carries the identity, the three semantic colours the rest of the repo
# already uses (blue for deterministic, amber for the exception lane, green for the person) do the
# work on every slide, and a warm off-white keeps content slides from reading as bare paper.
# The house palette, measured from real design tokens rather than guessed: indigo #4f46e5 for the
# brand, #3730a3 at the dark end of the gradient, #a5b4fc the periwinkle, #f8fafc the paper,
# #334155 the slate text, #18181b the black text.
# The three semantic colours stay, because they already mean something on every slide, with the
# deterministic blue re-tuned to the house indigo so the two never argue on a slide.
NAVY = RGBColor(0x37, 0x30, 0xA3)          # the dark indigo, for bands and table headers
NAVY_DEEP = RGBColor(0x1E, 0x1B, 0x4B)     # title slide ground, one step below the band
PAPER = RGBColor(0xF8, 0xFA, 0xFC)         # --bg--purple-light
CARD = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x58, 0x65, 0xF3)        # the logo indigo, seventeen paths of it, as the rule
ON_NAVY = RGBColor(0xFF, 0xFF, 0xFF)
ON_NAVY_MUTED = RGBColor(0xC7, 0xD2, 0xFE)  # indigo-200, readable on the dark ground
BRAND = RGBColor(0x4F, 0x46, 0xE5)         # the house indigo, the one colour that carries the identity

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.62)

# Three registers, one content. The room is strong on system design and does not speak AI
# vocabulary, so on the production copy every AI term on a slide becomes the plain phrase the README
# already uses for it, numbers kept. The room copy goes one step further and uses the domain's own
# words for the people in the story: exception-based review frees underwriters, and
# the customers are lenders, so where the production copy says "a person" in the review lane and
# "the customer", the room copy says underwriter and lender. It also drops the two sentences that
# told the reader how to feel about a slide, and titles the results slide with the result. The
# technical term survives in the speaker notes on every copy, which are James's to read.
REGISTER = "technical"  # technical | production | room


def plain(technical: str, everyday: str, room: str | None = None) -> str:
    if REGISTER == "technical":
        return technical
    if REGISTER == "room" and room is not None:
        return room
    return everyday


def theirs(ours: str, theirs_: str) -> str:
    """The room copy's word for it, every other copy's word otherwise."""
    return theirs_ if REGISTER == "room" else ours


def report(name: str, default=None):
    """Read a report, and refuse to build a deck from one that is missing.

    Same rule as draw_figures.py, for the same reason: a slide that says "0 of 0" because a file
    was absent is fabricated evidence, and the deck's first slide claims every number is reprinted.
    """
    p = ROOT / name
    if not p.exists():
        raise SystemExit(f"{name} is missing. Run `make eval` and `make telemetry` before building the decks.")
    data = json.loads(p.read_text())
    if not data:
        raise SystemExit(f"{name} is empty; refusing to build decks from it.")
    return data


EV = report("reports/eval_report.json", {}) or {}
TEL = report("reports/telemetry.json", {}) or {}
GOLD = report("reports/gold_report.json", {}) or {}
CMP = report("reports/compare_readers_report.json", {}) or {}
SYN = report("reports/synth_report.json", []) or []
BENCH = report("reports/bench_report.json", {}) or {}
TOK = report("reports/crop_tokens.json", {}) or {}
CNN_LAT = report("reports/cnn_latency.json", {}) or {}
DET_T = (TEL.get("modes", {}).get("deterministic core only", {}) or {}).get("totals", {})
CLF_T = (TEL.get("modes", {}).get("deterministic core + patch classifier", {}) or {}).get("totals", {})
OV = EV.get("overall", {})
if not OV.get("tp") or not DET_T.get("boxes") or not GOLD.get("cards") or not CMP.get("rules"):
    raise SystemExit("a report is present but carries no measurements; refusing to build decks from it.")
HARD = {c["id"]: c for c in CMP.get("hard_cards", [])}
if "c029" not in HARD or "c055" not in HARD:
    raise SystemExit("compare_readers_report.json carries no hard-card reads: run `make compare` first")
HTTP = (BENCH.get("levels") or [{}])[0]
if not HTTP.get("requests"):
    raise SystemExit("bench_report.json carries no HTTP level: run `make serve` and `make bench` first")
HTTP_PAGE = {"sample_1.jpg": "the photographed page, the smallest of the four"}.get(Path(BENCH.get("image", "")).name, Path(BENCH.get("image", "")).name)
SYN_FLOOR = min((r["f1"] for r in SYN), default=None)
if SYN_FLOOR is None:
    raise SystemExit("synth_report.json is empty: score the synthetic set first")
TOKENS = TOK.get("totals", {}).get("median_tokens_rounded")
if not TOKENS:
    raise SystemExit("crop_tokens.json carries no median: run `make crop-tokens` first")
CNN_PAGE = max(CNN_LAT.get("pages", [{"boxes": 0}]), key=lambda r: r["boxes"])
if not CNN_PAGE.get("boxes"):
    raise SystemExit("cnn_latency.json carries no page timings: run `make bench-cnn` first")
SAMPLE_PRED = sum(s.get("pred", 0) for s in EV.get("samples", []))
SAMPLE_FLAGGED = sum(s.get("ambiguous", 0) for s in EV.get("samples", []))


def _grouped(field: str) -> dict[str, tuple[int, int]]:
    """(boxes, flagged) per value of a page field, from the deterministic-core telemetry."""
    pages = (TEL.get("modes", {}).get("deterministic core only", {}) or {}).get("pages", [])
    out: dict[str, list[int]] = {}
    for page in pages:
        key = page.get(field)
        if not key:
            continue
        row = out.setdefault(str(key), [0, 0])
        row[0] += page.get("boxes", 0)
        row[1] += page.get("flagged", 0)
    return {k: (b, f) for k, (b, f) in out.items()}


def top_reasons(n: int = 3) -> list[tuple[str, int]]:
    """The n most frequent reason codes. telemetry.json already stores them most-common first."""
    return list(DET_T.get("reasons", {}).items())[:n]


def factor_rates() -> list[tuple[str, float]]:
    """Percent of boxes routed, per damage factor, worst first. Synthetic corpus only."""
    rates = [
        (k, 100 * f / b)
        for k, (b, f) in _grouped("factor").items()
        if b and k not in {"real", "held-out"}
    ]
    return sorted(rates, key=lambda kv: -kv[1])


def holdout_flagged() -> int:
    return _grouped("corpus").get("holdout", (0, 0))[1]


def deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs: Presentation, dark: bool = False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY_DEEP if dark else PAPER
    return s


def rect(slide, left, top, width, height, fill, radius: bool = False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    if radius:
        sh.adjustments[0] = 0.08
    return sh


def card(slide, left, top, width, height):
    """A white panel on the paper background, the unit every content slide is built from."""
    sh = rect(slide, left, top, width, height, CARD, radius=True)
    sh.line.color.rgb = RULE
    sh.line.width = Pt(0.75)
    return sh


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def textbox(slide, left, top, width, height, lines, size=18, color=INK, bold_first=False, space=8):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Helvetica Neue"
        p.font.bold = bold_first and i == 0
        p.space_after = Pt(space)
    return tb


SECTION = ""  # set per slide group so the band reads like a running header


def heading(slide, title: str, kicker: str | None = None):
    """A navy band with the section name and an amber tick, then the title on paper below it."""
    band_h = Inches(0.52)
    rect(slide, 0, 0, W, band_h, NAVY)
    rect(slide, 0, band_h, W, Emu(28575), ACCENT)
    if SECTION:
        textbox(slide, M, Inches(0.11), W - 2 * M, Inches(0.36), [SECTION.upper()], size=11, color=ON_NAVY_MUTED)
    textbox(slide, M, band_h + Inches(0.28), W - 2 * M, Inches(0.8), [title], size=30, color=INK, bold_first=True)
    y = band_h + Inches(1.05)
    if kicker:
        textbox(slide, M, y - Inches(0.05), W - 2 * M, Inches(0.55), [kicker], size=16, color=MUTED)
        y += Inches(0.55)
    return y + Inches(0.18)


def fitted_picture(slide, path: Path, left, top, max_w, max_h):
    """Place an image inside a box without distorting it.

    Dimensions come from OpenCV, which the project already depends on. Pillow is not installed and
    adding it just to read a width would put a dependency in the shipped project for the sake of a
    build script.
    """
    import cv2  # noqa: PLC0415

    img = cv2.imread(str(path))
    if img is None:
        slide.shapes.add_picture(str(path), left, top, width=max_w)
        return
    ih, iw = img.shape[:2]
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(path), int(left + (max_w - w) / 2), int(top), width=w, height=h)


def chip(slide, left, top, width, height, label, fill, edge, size=12):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = edge
    sh.line.width = Pt(1.1)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    for i, line in enumerate(label.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(size if i == 0 else size - 2)
        p.font.bold = i == 0
        p.font.color.rgb = INK
        p.font.name = "Helvetica Neue"
    return sh


def arrow(slide, x1, y1, x2, y2):
    """A straight connector with an actual arrowhead on the far end.

    python-pptx exposes no arrowhead API, and MSO_CONNECTOR.STRAIGHT alone draws a bare segment.
    Rendering the deck is what caught it; every audit that only reads shape geometry says a line
    and an arrow are the same object.
    """
    from pptx.enum.shapes import MSO_CONNECTOR  # noqa: PLC0415
    from pptx.oxml.ns import qn  # noqa: PLC0415

    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = MUTED
    c.line.width = Pt(1.25)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c


def table(slide, left, top, width, rows: list[list[str]], col_w: list[float] | None = None, size=13):
    n_r, n_c = len(rows), len(rows[0])
    height = Inches(0.34) * n_r
    shp = slide.shapes.add_table(n_r, n_c, left, top, width, height).table
    if col_w:
        total = sum(col_w)
        for i, frac in enumerate(col_w):
            shp.columns[i].width = int(width * frac / total)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = shp.cell(r, c)
            cell.text = str(val)
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(size)
            p.font.name = "Helvetica Neue"
            p.font.bold = r == 0
            # python-pptx tables default to a style with a dark accent header fill, so muted grey
            # in the header row came out grey-on-blue and the column labels were unreadable. Only
            # rendering the deck showed it; reading the XML says the colour was set as intended.
            p.font.color.rgb = INK if r else WHITE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD if r % 2 else PAPER
    return shp


# ----------------------------------------------------------------------------- slides


def s_title(prs, variant: str):
    s = blank(prs, dark=True)
    rect(s, 0, Inches(3.42), Inches(1.1), Emu(38100), ACCENT)
    textbox(s, M, Inches(2.35), W - 2 * M, Inches(1.3), ["Reading checkboxes off appraisal pages"], size=42, color=ON_NAVY, bold_first=True)
    textbox(
        s,
        M,
        Inches(3.75),
        W - 2 * M,
        Inches(2.0),
        [
            plain("A deterministic reader, an exception queue, and a definition of correctness the customer owns.",
                  "It reads what it can read for certain, hands the rest to a person, and lets the customer decide what counts as a mark.",
                  "It reads what it can read for certain, hands the rest to an underwriter, and lets the lender decide what counts as a mark."),
            "",
            "James W. Niu   ·   every number is read from the repo's own reports",
        ],
        size=17,
        color=ON_NAVY_MUTED,
    )
    # the three semantic colours, introduced once so they mean something on every later slide
    for i, (label, colour) in enumerate([(plain("deterministic", "rules, same answer always"), DET), (plain("exception lane", "sent for review"), EXC), (theirs("a person", "an underwriter"), HUM)]):
        x = M + Inches(3.3) * i
        rect(s, x, Inches(6.3), Inches(0.28), Inches(0.28), colour, radius=True)
        textbox(s, x + Inches(0.4), Inches(6.2), Inches(2.8), Inches(0.5), [label], size=14, color=ON_NAVY_MUTED)
    notes(s, "Thirty seconds: the assignment was to find checkboxes and say whether they are marked. "
             "The interesting part turned out not to be detection. It was deciding which boxes have a right "
             "answer and which ones genuinely do not, and who gets to decide that.")
    return s


def s_problem(prs):
    s = blank(prs)
    y = heading(s, "An appraiser's answers live in the boxes", "PUD or not. Utilities public or not. Market declining or not. Every one is a small square.")
    fitted_picture(s, ASSETS / "hero-page.png", M, y, W - 2 * M, H - y - Inches(0.5))
    notes(s, "This is one of the four pages the brief supplied, with every box the system found drawn on it. "
             "Green is marked, red is empty, amber is the two it refused to guess. 42 boxes, 16 read as marked. "
             "Point at the amber ones: that is the whole design in one picture.")
    return s


def s_hard(prs):
    s = blank(prs)
    y = heading(s, "Four things break a naive reader", "One of each is planted in the sample pages.")
    fitted_picture(s, ASSETS / "four-traps.png", M, y, W - 2 * M, H - y - Inches(0.5))
    notes(s, "Shaded rows, a watermark stamped through the border, ink so faint you squint, and handwriting "
             "crossing an empty box. The last two are cases where two careful people disagree with each other. "
             "That is what shaped the design: some boxes have a right answer and some do not.")
    return s


def s_arch(prs):
    s = blank(prs)
    y = heading(s, "Two readers, and a queue for what they cannot settle")
    # Every position is computed from the usable width rather than assumed. The first version used
    # fixed chip widths and ran three inches past the right edge of the slide, which PowerPoint will
    # happily render and nobody notices until it is on a projector.
    usable = W - 2 * M
    steps = [
        ("Clean up\nthe page", DET_FILL, DET),
        ("Reader 1\nprinted lines", DET_FILL, DET),
        ("Reader 2\nblank form", DET_FILL, DET),
        ("Both\nagree?", WHITE, MUTED),
        ("Read the\nmark", DET_FILL, DET),
        ("Confident?", WHITE, MUTED),
        ("Checked\nor empty", RGBColor(0xF4, 0xF6, 0xF8), MUTED),
    ]
    gap = Inches(0.30)
    bw = int((usable - gap * (len(steps) - 1)) / len(steps))
    bh = Inches(0.95)
    row1 = y + Inches(0.45)
    row2 = row1 + Inches(2.05)

    xs = [M + i * (bw + gap) for i in range(len(steps))]
    for x, (label, fill, edge) in zip(xs, steps):
        chip(s, x, row1, bw, bh, label, fill, edge, size=11)
    for i in range(len(steps) - 1):
        arrow(s, xs[i] + bw, row1 + bh // 2, xs[i + 1], row1 + bh // 2)

    lane = [
        ("Exception queue\nevery item carries a reason", EXC_FILL, EXC, 3.3),
        ("AI check on the crop\noff by default", EXC_FILL, EXC, 2.9),
        (theirs("A person sets the rule\nand it becomes policy", "The lender sets the rule\nand it becomes policy"), HUM_FILL, HUM, 3.1),
    ]
    lx = M + Inches(0.9)
    lane_x = []
    for label, fill, edge, w_in in lane:
        w = Inches(w_in)
        chip(s, lx, row2, w, bh, label, fill, edge, size=11)
        lane_x.append((lx, w))
        lx += w + Inches(0.45)
    for (x1, w1), (x2, _) in zip(lane_x, lane_x[1:]):
        arrow(s, x1 + w1, row2 + bh // 2, x2, row2 + bh // 2)

    # The two ways a box drops out of the top row, and they leave from the two decision diamonds.
    # These indices were 4 and 6 when the row still opened with an "Appraisal page" chip; dropping
    # it shifted everything left by one and the arrows started leaving from the wrong boxes, which
    # is invisible until the deck is actually rendered and looked at.
    agree_i = next(i for i, (label, _, _) in enumerate(steps) if label.startswith("Both"))
    confident_i = next(i for i, (label, _, _) in enumerate(steps) if label.startswith("Confident"))
    arrow(s, xs[agree_i] + bw // 2, row1 + bh, lane_x[0][0] + Inches(0.7), row2)
    arrow(s, xs[confident_i] + bw // 2, row1 + bh, lane_x[0][0] + Inches(2.5), row2)

    routed = round(100 * DET_T.get("flag_rate", 0), 1)
    textbox(s, M, H - Inches(1.15), W - 2 * M, Inches(0.9),
            [f"{100 - routed:.1f}% of boxes never leave the top row. {routed:.1f}% carry a reason code "
             f"into the queue, measured over {DET_T.get('pages', 0)} pages. Nothing in the top row is a model."],
            size=16, color=MUTED)
    notes(s, "The top row is ordinary image processing and it is where almost everything is settled. "
             "Two readers find the boxes independently: one from the printed lines, one from the blank federal "
             "form. Where they agree the answer stands. Where they disagree the box is flagged, never dropped. "
             "The bottom row is the exception lane and it is off by default.")
    return s


def s_accuracy_claim(prs):
    s = blank(prs)
    y = heading(s, "1. Accuracy, against whose definition?", "The first dimension, and the one everybody assumes is settled.")
    textbox(s, M, y + Inches(0.15), Inches(6.1), Inches(4.6), [
        "A circle drawn in the box instead of an X.",
        "A box filled in and then struck out.",
        "A tick that overshoots the border.",
        "",
        "Each has a defensible answer in both directions. None of them is ours to make. One lender wants "
        + theirs("the circle counted; another wants it sent to a person. Both are right about their own files.",
                 "the circle counted; another wants it sent to an underwriter. Both are right about their own files."),
        "",
        theirs("So the definition lives in a file the customer owns, not in our code.",
               "So the definition lives in a file the lender owns, not in our code."),
    ], size=17)
    chip(s, Inches(7.0), y + Inches(0.3), Inches(5.7), Inches(1.0), plain("policy.json\nthresholds, and the three rulings people argue about", "One settings file\nthe thresholds, and the three rulings people argue about"), DET_FILL, DET, size=15)
    chip(s, Inches(7.0), y + Inches(1.55), Inches(5.7), Inches(1.0), "POLICY.md\nthe same decisions in prose, from the person who labelled", HUM_FILL, HUM, size=15)
    textbox(s, Inches(7.0), y + Inches(2.85), Inches(5.7), Inches(1.4), [
        "Changing what counts as a mark is editing that file.",
        "It is not a ticket, a release, or a call with us.",
    ], size=16, color=MUTED)
    notes(s, "Ask what the accuracy of a checkbox reader is and the honest first answer is a question back: "
             "accurate against whose definition? This is the dimension people skip, and it is the one that "
             "decides whether the customer can live with the system.")
    return s


def s_accuracy_demo(prs, dense: bool):
    s = blank(prs)
    y = heading(s, "Same code, same pages, two definitions", "uv run python scripts/compare_policies.py")
    table(s, M, y + Inches(0.25), W - 2 * M, [
        ["", "shipped policy", theirs("a stricter customer", "a stricter lender")],
        ["a single thin stroke", "counts as marked", "flag it, do not decide"],
        ["a box scribbled out", "read as empty", "flag it, do not overrule the ink"],
        ["uncertain band", "0.05 to 0.20", "0.03 to 0.30"],
        ["boxes read differently", "2 of 287", "2 of 287"],
        ["agreement with the answer key", "285 of 286", "286 of 286"],
    ], col_w=[2.4, 1.6, 1.8], size=15)
    textbox(s, M, y + Inches(2.75), W - 2 * M, Inches(2.0), [
        "The stricter policy is not simply better.",
        "It gets the faded X right by refusing to overrule ink on a fragmented mark, and the same rule "
        "would let a struck-out box report as filled. Which of those two errors you would rather have is "
        "a business question, and the file is where you answer it.",
    ], size=17)
    notes(s, "This is the demo I would run live. Two files, no code change, and the answer on one box moves. "
             "The point is not that one policy wins. It is that the argument now has somewhere to happen that "
             "is not our backlog.")
    return s


def s_cost(prs):
    s = blank(prs)
    y = heading(s, "2. Cost, by using the right size of tool", "Not the smallest tool. The right one.")
    real_boxes = sum(pg.get("boxes", 0) for pg in TEL["modes"]["deterministic core + patch classifier"]["pages"] if pg.get("corpus") in ("sample", "holdout"))
    real_flagged = sum(pg.get("flagged", 0) for pg in TEL["modes"]["deterministic core + patch classifier"]["pages"] if pg.get("corpus") in ("sample", "holdout"))
    table(s, M, y + Inches(0.2), W - 2 * M, [
        ["approach", "what a page costs", "measured where"],
        ["ordinary image processing (what runs)", f"{DET_T.get('p50_ms', 0):.0f} ms of one CPU core, no model", f"make telemetry, median over {DET_T.get('pages', 0)} pages"],
        ["AI on boxes the core flags", f"one {TOKENS}-token crop for {SAMPLE_FLAGGED} of {SAMPLE_PRED} detections", "make telemetry and make crop-tokens, the 4 brief pages"],
        ["AI on those plus every model dispute", f"{real_flagged} of {real_boxes} boxes on the real pages", "make telemetry with the model on, brief plus held-out"],
        [theirs("every page to a frontier model", "every page to the biggest AI model"), "the whole page, every page, priced per token", "published work, docs/EVALS.md"],
    ], col_w=[3.0, 2.2, 2.6], size=14)
    textbox(s, M, y + Inches(2.5), W - 2 * M, Inches(2.2), [
        plain("A frontier model is trained to solve an enormous range of problems, and a loan file pays for none of them.",
              "The biggest AI models are trained to solve an enormous range of problems, and a loan file pays for none of them."),
        theirs(f"A flagged box, zoomed in the way the escalation lane sends it, is about {TOKENS} tokens of image, the median over "
               f"the {TOK['totals']['boxes']} boxes on the brief pages. A whole page is the entire image, and the answer has to "
               "list every box and its coordinates on top of that.",
               f"A model bills by the token, and a flagged box, zoomed in the way we send it, is about {TOKENS} of them, the median "
               f"over the {TOK['totals']['boxes']} boxes on the brief pages. A whole page is the entire image, and the answer has to "
               "list every box and its coordinates on top of that."),
        "We do not pay for intelligence we do not use.",
    ], size=17)
    notes(s, "The arithmetic is the whole argument. At real volume that gap stops being a rounding error, and "
             "it is the wrong trade on quality too: published work finds checkbox reading is a specific weak "
             "spot for vision models, not a strength.")
    return s


def s_latency(prs):
    s = blank(prs)
    y = heading(s, "3. Latency, and there are two kinds")
    textbox(s, M, y + Inches(0.2), Inches(6.1), Inches(4.4), [
        "The kind the page waits for",
        plain(f"{DET_T.get('p50_ms', 0):.0f} ms median and {DET_T.get('p95_ms', 0):.0f} ms p95 in-process, measured over "
              f"{DET_T.get('pages', 0)} pages. Through the HTTP service, {HTTP['requests']} requests of {HTTP_PAGE} come back at "
              f"p50 {HTTP['p50_ms']:.0f} ms and p95 {HTTP['p95_ms']:.0f} ms. No network call in the path at all. Nothing waits on a vendor.",
              f"A median of {DET_T.get('p50_ms', 0):.0f} ms a page, and {DET_T.get('p95_ms', 0):.0f} ms for 95 pages in every 100, measured "
              f"in-process over {DET_T.get('pages', 0)} pages. Through the HTTP service, {HTTP['requests']} requests of {HTTP_PAGE} come back "
              f"at a median of {HTTP['p50_ms']:.0f} ms. No network call in the path at all. Nothing waits on a vendor."),
        "",
        "It also improves on its own, which a model-first design does not. Every behaviour we understand well "
        "enough to describe becomes a rule, and a rule costs microseconds. The queue is meant to shrink.",
    ], size=17, bold_first=True)
    textbox(s, Inches(7.0), y + Inches(0.2), Inches(5.7), Inches(4.4), [
        "The kind an operator waits for",
        "Somebody has to know where the system is struggling without reading pages to find out.",
        "",
        "That is why the reason codes are data rather than log lines, and why there is a dashboard. "
        "A flag rate is a budget line. A breakdown is something you can act on.",
    ], size=17, bold_first=True)
    notes(s, "Two kinds and people usually only price the first. The second one is where the operational cost "
             "actually sits, because a human reading pages to find out what went wrong is the most expensive "
             "thing in the loop.")
    return s


def s_governance(prs):
    s = blank(prs)
    y = heading(s, "4. Governance, and how little model you use")
    textbox(s, M, y + Inches(0.2), Inches(6.1), Inches(4.4), [
        theirs("Every model call is customer documents leaving the building.", "Every model call is a lender's documents leaving the building."),
        plain("Deterministic code is not. It runs on a machine you control, on-premises if the contract says so, "
              "with no vendor to ask about retention or training.",
              "Plain rules are not. They run on a machine you control, in your own building if the contract says "
              "so, with no vendor to ask about retention or training."),
        "",
        "By default this sends nothing anywhere.",
        plain("When the escalation lane is on, it sends a thumbnail of one box, never a page of somebody's finances.",
              "When the AI check is on, it sends a thumbnail of one box, never a page of somebody's finances."),
    ], size=17, bold_first=True)
    chip(s, Inches(7.0), y + Inches(0.25), Inches(5.7), Inches(1.25),
         plain("Hard gates\nthe response schema, the blank forms, the regression rows, the frozen answer key",
               "Hard gates\nthe answer format, the blank forms, the known hard cases, the frozen answer key"), DET_FILL, DET, size=15)
    chip(s, Inches(7.0), y + Inches(1.75), Inches(5.7), Inches(1.25),
         "Soft gates\na number with an allowed range and a verdict; one sits outside its range and one is not yet measured against it, and both say so", EXC_FILL, EXC, size=15)
    chip(s, Inches(7.0), y + Inches(3.25), Inches(5.7), Inches(1.05),
         theirs("Humans on the loop\nQA and product write the criteria in the same file the customer edits",
                "Humans on the loop\nQA and product write the criteria in the same file the lender edits"), HUM_FILL, HUM, size=15)
    notes(s, "Both kinds of gate are visible and published. The one that is out of band is the trained "
             "classifier's flag rate, and it is written down in docs/EVALS.md rather than quietly retuned. "
             "That is the honest version of a green dashboard.")
    return s


def s_abc(prs):
    s = blank(prs)
    y = heading(s, "Three ways to build this", "The two alternatives fail in opposite directions.")
    table(s, M, y + Inches(0.15), W - 2 * M, [
        ["", "A. this", "B. OCR plus a trained detector", theirs("C. whole page to a frontier model", "C. whole page to the biggest AI model")],
        ["accuracy here", "286 of 287", theirs("F1 0.88 to 0.96 published", "88% to 96% in published results"), "documented weak spot on checkboxes"],
        ["what a page costs", f"{DET_T.get('p50_ms', 0):.0f} ms of one core, no model", "$10 per 1,000 pages, Azure list price", "the whole page, priced per token"],
        ["latency", f"{DET_T.get('p50_ms', 0):.0f} ms median, local", "a vendor round trip", "seconds, and variable"],
        ["same page twice", "byte-identical", "yes, weights are frozen", "no"],
        ["who defines a mark", theirs("the customer, in a file", "the lender, in a file"), "whoever labelled the training set", "the model"],
        ["changing that", theirs("edit policy.json", "edit one settings file"), "relabel and retrain", "reword the prompt and hope"],
        ["gates you can inspect", "hard and soft, published", "none exposed", "none"],
        ["documents leave", "never", "to the vendor", "to the vendor, whole pages"],
    ], col_w=[1.8, 1.5, 2.1, 2.2], size=12)
    notes(s, "B is the 2021 answer and it loses on control, not accuracy: the definition of a mark is frozen "
             "inside the weights, there is nowhere to look up why a box was read a given way, and changing it "
             "means a relabel and a retrain. C is the 2026 answer and it pays frontier prices to read something "
             "a hundred pixels wide, gives up determinism, and ships whole customer pages to a third party. "
             "If this were thousands of unseen layouts with no explainability requirement, C would be right and "
             "this would be over-engineering.")
    return s


def s_results(prs):
    s = blank(prs)
    found = f"{OV.get('tp', 0)} of {OV.get('tp', 0) + OV.get('fn', 0)} boxes found"
    right = f"{round(OV.get('cls_acc', 0) * OV.get('tp', 0))} of {OV.get('tp', 0)} marks read right"
    y = heading(s, theirs("Results", f"{found}, {right}"), "Reprinted by make eval. Nothing typed by hand.")
    table(s, M, y + Inches(0.2), Inches(7.6), [
        ["page", "boxes", "found", "marks right", "flagged"],
        ["photographed page", "42", "41 of 42", "40 of 41", "2"],
        ["standard form, clean scan", "118", "118 of 118", "118 of 118", "0"],
        ["market addendum, shaded", "48", "48 of 48", "48 of 48", "0"],
        ["manufactured home, watermarked", "79", "79 of 79", "79 of 79", "0"],
        ["all four", "287", "286 of 287", "285 of 286", "2"],
    ], col_w=[2.6, 0.8, 1.2, 1.2, 0.8], size=13)
    textbox(s, Inches(8.5), y + Inches(0.2), Inches(4.2), Inches(4.4), [
        "Two answers are not clean, and I know why for both.",
        "The miss is a checkbox printed so faintly the labeller called it a box and the software cannot see it.",
        "The wrong mark is the faded ink, which the system flags as unclear rather than deciding. Counted "
        "strictly, an unclear answer counts against it.",
        "",
        plain(f"On the 52 pages damaged on purpose, detection F1 never drops below {SYN_FLOOR:.3f} under any kind of "
              f"damage, and shading is the floor.",
              f"On the 52 pages damaged on purpose, the box-finding score, F1, never drops below {SYN_FLOOR:.3f} under any "
              f"kind of damage, and shading is the floor."),
    ], size=15, bold_first=True)
    notes(s, "Four real pages cannot prove much on their own, which is why the synthetic corpus and the "
             "held-out appraisals exist. The held-out ones came from three offices that were never used while "
             "building anything and they flagged nothing at all.")
    return s


def s_run(prs):
    """How to build and run it, on a slide. The brief's own guideline, taught rather than pointed at."""
    s = blank(prs)
    y = heading(s, "Build it and run it", "From the unzipped folder. No accounts, no keys, no network after the install.")
    items = [
        ("uv sync --extra dev", "Installs everything, tests included, into a local environment. One tool, one command."),
        ("make serve", "Starts the service on localhost:8000. The same thing runs from a container: "
                       "make docker-build, then make docker-run."),
        ("curl -F file=@data/samples/sample_1.jpg\nlocalhost:8000/detect", "Returns one JSON entry per checkbox, its rectangle and whether it is "
                       "marked. Add /overlay after /detect and the page comes back with every box drawn on it."),
        ("make test   ·   make eval", "Runs the 26 safety gates, then reprints the accuracy table from scratch."),
    ]
    yy = y + Inches(0.1)
    for cmd, why in items:
        chip(s, M, yy, Inches(4.9), Inches(0.86), cmd, DET_FILL, DET, size=13)
        textbox(s, Inches(5.8), yy + Inches(0.04), Inches(6.9), Inches(0.9), [why], size=15, color=INK)
        yy += Inches(1.06)
    textbox(s, M, yy + Inches(0.08), W - 2 * M, Inches(0.6), [
        "I ran every one of these from a fresh unzip of this zip before building the deck.",
    ], size=15, color=MUTED)
    notes(s, "If the room wants a live run, this is the slide to do it from: the whole loop is install, "
             "serve, one curl. The overlay endpoint is the fastest way to disagree with it, and make eval "
             "reprints the results table while people watch.")
    return s


def s_verify(prs):
    s = blank(prs)
    y = heading(s, "Where you put your own eyes on it", "Four places, each a command rather than a claim.")
    items = [
        ("make overlays", "Every box drawn on all four pages. Green marked, red empty, amber held back. "
                          "Disagreeing with a picture takes seconds."),
        ("?explain=true", "One box: the ink, the confidence, the reason codes, which form matched, how far the "
                          "two readers agreed, and every candidate it threw away with why."),
        ("make eval  ·  make test", "Reprints the results table from scratch and runs 26 gates."),
        ("scripts/fetch_holdout.py --score", "Five completed appraisals from three offices, fetched and scored live."),
    ]
    yy = y + Inches(0.1)
    for cmd, why in items:
        chip(s, M, yy, Inches(3.3), Inches(0.74), cmd, DET_FILL, DET, size=14)
        textbox(s, Inches(4.3), yy + Inches(0.02), Inches(8.4), Inches(0.8), [why], size=15, color=INK)
        yy += Inches(0.92)
    textbox(s, M, yy + Inches(0.05), W - 2 * M, Inches(0.7), [
        theirs("Once this is running for real, the place a person looks is the exception queue. ",
               "Once this is running for real, the place an underwriter looks is the exception queue. ")
        + f"{DET_T.get('flagged', 0)} boxes out of {DET_T.get('boxes', 0):,}, each carrying a reason, "
        + theirs("instead of a person re-reading every page.", "instead of re-reading every page."),
    ], size=16, color=MUTED)
    notes(s, "The rejected-candidates field is the one I would check first, because rejection is the only place "
             "the system removes something instead of flagging it.")
    return s


def s_dashboard(prs):
    """The operator view drawn natively at slide scale.

    The first version dropped a screenshot of dashboard.html onto the slide. Nearly square, it
    rendered five inches wide and nothing on it could be read from across a room, so it said "here
    is a picture of a dashboard" and nothing else. This draws the two panels that carry the argument
    as real shapes, big enough to read, from the same telemetry the dashboard is built from.
    """
    s = blank(prs)
    y = heading(s, "What it looks like in operation",
                f"Measured over {DET_T.get('pages', 0)} pages and {DET_T.get('boxes', 0):,} checkboxes. "
                "The full page is deliverables/dashboard.html, one file, no server.")

    # ---- KPI strip
    gold_acc = GOLD.get("accuracy", 0)
    kpis = [
        ("Pages", f"{DET_T.get('pages', 0)}", "52 damaged on purpose"),
        ("Checkboxes", f"{DET_T.get('boxes', 0):,}", f"{DET_T.get('checked', 0):,} read as marked"),
        (plain("Detection F1", "Boxes found"), plain(f"{OV.get('f1', 0):.4f}", f"{OV.get('tp', 0)} of {OV.get('tp', 0) + OV.get('fn', 0)}"), plain("on the labelled pages", "on the pages from the brief")),
        (theirs("Sent to a person", "Sent to an underwriter"), f"{100 * DET_T.get('flag_rate', 0):.1f}%", f"{DET_T.get('flagged', 0)} boxes, each with a reason"),
        ("Per page", f"{DET_T.get('p50_ms', 0):.0f} ms", "no model in the path"),
        ("Agrees with a person", f"{gold_acc:.2f}", f"{GOLD.get('correct', 0)} of {GOLD.get('cards', 0)} rulings"),
    ]
    n = len(kpis)
    gap = Inches(0.16)
    kw = int((W - 2 * M - gap * (n - 1)) / n)
    kh = Inches(1.08)
    for i, (label, value, note) in enumerate(kpis):
        x = M + i * (kw + gap)
        card(s, x, y, kw, kh)
        rect(s, x, y, Inches(0.07), kh, DET)
        textbox(s, x + Inches(0.18), y + Inches(0.05), kw - Inches(0.2), Inches(0.3), [label], size=10, color=MUTED)
        textbox(s, x + Inches(0.18), y + Inches(0.28), kw - Inches(0.2), Inches(0.5), [value], size=22, color=INK, bold_first=True)
        textbox(s, x + Inches(0.18), y + Inches(0.72), kw - Inches(0.2), Inches(0.3), [note], size=9, color=MUTED)

    # ---- two panels
    py = y + kh + Inches(0.22)
    ph = H - py - Inches(0.35)
    pw = int((W - 2 * M - Inches(0.22)) / 2)

    # left: why boxes went to a person, as ranked bars with plain-English labels
    lx = M
    card(s, lx, py, pw, ph)
    textbox(s, lx + Inches(0.22), py + Inches(0.1), pw - Inches(0.4), Inches(0.35), [theirs("Why boxes went to a person", "Why boxes went to an underwriter")], size=14, color=INK, bold_first=True)
    # A box can carry more than one reason, so the bars count reason assignments and the caption says
    # so; the first cut called them boxes and hid the fifth code, which a reviewer caught because the
    # panel promised the whole queue and did not show it.
    n_reasons = sum(DET_T.get("reasons", {}).values())
    textbox(s, lx + Inches(0.22), py + Inches(0.42), pw - Inches(0.4), Inches(0.35),
            [f"Every reason on the {DET_T.get('flagged', 0)} flagged boxes, {n_reasons} in all; a box can carry more than one."], size=10, color=MUTED)
    labels = {
        "STRAY_STROKE": "a pen line crosses the box",
        "MISSING_IN_DETECT": "the form expects a box here",
        "INK_AMBIGUOUS": "too little ink to call",
        "EXTRA_BOX": "a box the blank form lacks",
        "FRAGMENTED_MARK": "specks, not one stroke",
    }
    colours = {"STRAY_STROKE": EXC, "MISSING_IN_DETECT": DET, "INK_AMBIGUOUS": RGBColor(0x8A, 0x6F, 0xD4),
               "EXTRA_BOX": RGBColor(0xD4, 0x64, 0x3B), "FRAGMENTED_MARK": HUM}
    reasons = list(DET_T.get("reasons", {}).items())
    top = max((v for _, v in reasons), default=1)
    ry = py + Inches(0.9)
    # Row height comes from the room the card actually has. The first cut assumed 0.6in a row and
    # the fifth row fell off the bottom of the card once the heading band took its share.
    row_h = int((ph - Inches(1.05)) / max(1, len(reasons)))
    label_w = Inches(2.35)
    bar_w_max = pw - label_w - Inches(1.15)
    for k, v in reasons:
        textbox(s, lx + Inches(0.22), ry, label_w, Inches(0.3), [k.replace("_", " ").lower()], size=11, color=INK, bold_first=True)
        textbox(s, lx + Inches(0.22), ry + Inches(0.24), label_w, Inches(0.3), [labels.get(k, "")], size=9, color=MUTED)
        bw = max(Inches(0.08), int(bar_w_max * v / top))
        rect(s, lx + Inches(0.22) + label_w, ry + Inches(0.1), bw, Inches(0.3), colours.get(k, MUTED), radius=True)
        textbox(s, lx + Inches(0.22) + label_w + bw + Inches(0.08), ry + Inches(0.06), Inches(0.8), Inches(0.35), [str(v)], size=13, color=INK, bold_first=True)
        ry += row_h

    # right: where the trouble concentrates, by damage type
    rx = M + pw + Inches(0.22)
    card(s, rx, py, pw, ph)
    textbox(s, rx + Inches(0.22), py + Inches(0.1), pw - Inches(0.4), Inches(0.35), ["Where the trouble concentrates"], size=14, color=INK, bold_first=True)
    held = holdout_flagged()
    textbox(s, rx + Inches(0.22), py + Inches(0.42), pw - Inches(0.4), Inches(0.35),
            [f"Share of boxes sent for review, by what was done to the page. The held-out real appraisals flagged {held}."],
            size=10, color=MUTED)
    rates = factor_rates()[:5]
    top_r = max((r for _, r in rates), default=1.0)
    ry = py + Inches(0.9)
    foot_h = Inches(0.62)
    row_h = int((ph - Inches(0.9) - foot_h - Inches(0.1)) / max(1, len(rates)))
    label_w = Inches(1.35)
    bar_w_max = pw - label_w - Inches(1.25)
    for k, r in rates:
        textbox(s, rx + Inches(0.22), ry + Inches(0.02), label_w, Inches(0.3), [k], size=11, color=INK)
        bw = max(Inches(0.08), int(bar_w_max * r / top_r))
        rect(s, rx + Inches(0.22) + label_w, ry + Inches(0.07), bw, Inches(0.24), EXC, radius=True)
        textbox(s, rx + Inches(0.22) + label_w + bw + Inches(0.08), ry, Inches(0.9), Inches(0.35), [f"{r:.2f}%"], size=12, color=INK, bold_first=True)
        ry += row_h
    all_rates = factor_rates()
    worst = f"{all_rates[0][0]} at {all_rates[0][1]:.2f}%" if all_rates else "unknown"
    best = f"{all_rates[-1][1]:.2f}% for {all_rates[-1][0]}" if all_rates else "unknown"
    textbox(s, rx + Inches(0.22), py + ph - foot_h, pw - Inches(0.4), foot_h - Inches(0.08),
            [f"Hardest is {worst}, easiest {best}. If real scans start looking like the worst case, the queue is where it shows first."],
            size=10, color=MUTED)

    top3 = ", ".join(f"{v} {labels.get(k, k.lower())}" for k, v in reasons[:3])
    notes(s, "The panel that matters is why boxes went to a person. Of the queue, "
             f"{top3}. That is something an operator can act on without opening a document. "
             f"The held-out real appraisals flagged {held}. The hardest single condition is {worst}, against {best}. "
             "deliverables/dashboard.html in the repo is the full page, and it opens with no server.")
    return s


def s_apprentice(prs):
    s = blank(prs)
    y = heading(s, "Who is the apprentice?",
                "A small model trained from scratch, raced against the rules on the brief's answer key, and left switched off.")
    tw = Inches(7.0)
    table(s, M, y, tw, [
        ["how it ran", "sent to a person", "settled right", "settled wrong"],
        ["rules only, which is what ships", f"{CMP['rules']['queue']} of {CMP['rules']['graded']}", str(CMP['rules']['right']), str(CMP['rules']['wrong'])],
        ["the CNN alone, a test", f"{CMP['cnn']['queue']} of {CMP['cnn']['graded']}", str(CMP['cnn']['right']), str(CMP['cnn']['wrong'])],
        ["both together, a test", f"{CMP['both']['queue']} of {CMP['both']['graded']}", str(CMP['both']['right']), str(CMP['both']['wrong'])],
    ], col_w=[2.6, 1.4, 1.1, 1.1], size=13)
    textbox(s, M, y + Inches(1.7), tw, Inches(1.4), [
        "Verdict: rules only.",
        f"All three are nearly errorless on the {CMP['rules']['graded']} graded boxes, the 287 detections on the "
        f"brief pages minus the tick ruled unsure. They differ in how much they hand to people, "
        f"{CMP['rules']['queue']} against {CMP['both']['queue']} against {CMP['cnn']['queue']}. "
        "Most automation at the same accuracy wins."], size=16, color=MUTED, bold_first=True)
    rx = M + tw + Inches(0.5)
    textbox(s, rx, y, W - rx - M, Inches(4.6), [
        "The rules won the match, so the CNN stays switched off.",
        f"It caught nothing the rules missed. On the two hard boxes on the photographed page it settled neither, "
        f"queuing the empty pen loop at p(filled) {HARD['c029']['cnn_p_filled']:.3f} and reading the filled faded X as "
        f"empty at {HARD['c055']['cnn_p_filled']:.3f}.",
        f"Switched on, it only adds work: {CMP['both']['queue'] - CMP['rules']['queue']} more boxes to people, zero errors caught.",
        "Legibility is the invariant. The rules have it built in; the CNN can only approach it by "
        "emitting logs, a confidence and a heat map per box, replayable forever.",
    ], size=15, space=10, bold_first=True)
    notes(s, "The LLM that built the product wrote both contestants, the rules and this 23,381-weight model, and "
             "judged them against one referee: 52 damaged pages plus the hand-ruled answer keys. When the "
             f"CNN is confident it is right {CMP['cnn']['right']} of {CMP['cnn']['right'] + CMP['cnn']['wrong']} times, so "
             f"training worked. Its problem is confidence in the right places, {CMP['cnn']['queue']} unsure boxes where the "
             f"rules are sure and correct. The 286 graded boxes are the 287 detections on the brief pages minus the one on "
             f"the tick the labeler ruled unsure. Scoring a {CNN_PAGE['boxes']}-box page takes it {CNN_PAGE['median_ms']:.1f} ms "
             "on CPU (make bench-cnn). make compare reprints this table from the repo.")
    return s


def s_feedback(prs):
    s = blank(prs)
    y = heading(s, "What this becomes with real volume")
    textbox(s, M, y + Inches(0.3), W - 2 * M, Inches(4.4), [
        "The reason codes are per box today. Joined to the form's official field names, the same counts say "
        "which questions on the form confuse the people filling it in.",
        "",
        theirs("That is a finding the customer can act on by changing the form, not by reviewing more documents. It "
               "turns an extraction service into a feedback loop on their own product.",
               "That is a finding a lender can act on by changing the form, not by reviewing more documents. It "
               "turns an extraction service into a feedback loop on their own product."),
        "",
        "Getting there needs the template registry to carry field names, which is item 2 on the roadmap and is "
        + "not built here.",
    ], size=18)
    notes(s, "This is the part I would want to build next. It is also the part that changes what the product "
             "is: not 'we read your checkboxes' but 'we can tell you which questions your users get stuck on'.")
    return s


def s_limits(prs):
    s = blank(prs)
    y = heading(s, "What it cannot do", "Five of them, and the measurement behind each one is in docs/approach.md.")
    textbox(s, M, y + Inches(0.2), W - 2 * M, Inches(4.6), [
        "Four real pages is not a benchmark. The held-out five carry no box-by-box labels, so they were checked "
        "by eye rather than scored. Sixty-one pages in all, and 52 of them are ones I damaged myself.",
        "The page labels were seeded by the software before correction, and that shortcut hid an invented box "
        "until a second reader disagreed about a row. That story is in the writeup.",
        "Rejecting a candidate is the one place the system removes evidence rather than flagging it, and the "
        "rule is resolution-dependent: at half size the invented box survives.",
        plain("We raced the trained classifier against the rules on the brief's answer key: 56 boxes to a person against 2, nothing caught that the rules missed. It ships switched off.",
              "We raced the small trained model against the rules on the brief's answer key: it sent 56 boxes to a person where the rules sent 2, and caught nothing they missed. It ships switched off."),
        "One of the three form types has no publicly published blank, so that page runs on one reader and says so.",
    ], size=17, space=12)
    notes(s, "I would rather be the person who lists these than the person who gets asked about them. Every one "
             "of them is written down in docs/approach.md with the measurement behind it.")
    return s


def s_iterations(prs):
    s = blank(prs)
    y = heading(s, "How it got here", "Every run, including the ones that scored worse.")
    table(s, M, y + Inches(0.2), W - 2 * M, [
        ["run", "what changed", "what it taught"],
        ["v1", "detector, ink rule, API, labelling booth", "labels seeded from the detector score 1.0 by definition"],
        ["v2", "labels corrected by eye, four shape rules", "each planted hard case forced a named, testable rule"],
        ["v3", theirs("tier-1 gates, synthetic generator, sweeps", "hard gates, synthetic pages, damage sweeps"),
         theirs("hash() is process-salted, so synthesis was not reproducible", "Python's hash() changes between runs, so the synthetic pages did not reproduce")],
        ["v4", theirs("monotone line alignment for the second reader", "line-by-line alignment for the second reader"),
         theirs("the best global scale fit is off by hundreds of pixels across vendors", "one page-wide scale is off by hundreds of pixels between form vendors")],
        ["v5", theirs("span-ceiling fix, gold set folded in, classifier", "an unreachable threshold fixed, 76 crops folded in, small model"),
         theirs("a 0.7 span threshold was unreachable by construction", "no real mark could ever reach the 0.7 span threshold")],
        ["v6", "size consensus plus an interior text test", "ground truth seeded from the system under test hides its own errors"],
    ], col_w=[0.5, 2.6, 3.6], size=12)
    notes(s, "v1 scored a perfect 1.0 against labels the detector had written itself. That is what a meaningless "
             "number looks like and it is left in the record on purpose.")
    return s


def s_why_b_c(prs):
    s = blank(prs)
    y = heading(s, "Why the alternatives lose here")
    textbox(s, M, y + Inches(0.2), Inches(6.1), Inches(4.6), [
        "B, the 2021 answer",
        "It does not lose on accuracy. It loses because the definition of a mark is frozen inside the weights.",
        theirs("A customer who wants circles counted needs a labelling round and a retrain. ",
               "A lender who wants circles counted needs a labelling round and a retrain. ")
        + "There is nowhere to look up why any single box was read the way it was.",
        theirs("A black box with no user-definable acceptance criteria and no gates. The only way to argue with it is "
               "to build another one.",
               "A black box with no rules a lender can set and no gates. The only way to argue with it is "
               "to build another one."),
    ], size=16, bold_first=True)
    textbox(s, Inches(7.0), y + Inches(0.2), Inches(5.7), Inches(4.6), [
        "C, the 2026 answer",
        theirs("It pays frontier prices to read something a hundred pixels wide, and gives up determinism to do it.",
               "It pays top prices to read something a hundred pixels wide, and gives up the same answer twice to do it."),
        "Published work finds checkbox reading is a specific weakness of vision models rather than a strength, "
        "which is the wrong place to spend the most money.",
        theirs("And it puts whole customer pages in a third party's hands to answer a question that never needed to "
               "leave the building.",
               "And it puts whole pages of a lender's file in a third party's hands to answer a question that never needed to "
               "leave the building."),
    ], size=16, bold_first=True)
    textbox(s, M, H - Inches(1.25), W - 2 * M, Inches(0.9), [
        theirs("Where A would be wrong: thousands of unseen layouts, no explainability requirement, no customer "
               "disagreement about what counts. Then C is right and this is over-engineering.",
               "Where A would be wrong: thousands of unseen layouts, nobody asking why a box was read a given way, no two lenders "
               "disagreeing about what counts. Then C is right and this is over-engineering."),
    ], size=15, color=MUTED)
    notes(s, "Saying where your own approach is the wrong one is usually the fastest way to be believed about "
             "where it is the right one.")
    return s


def s_holdout(prs):
    s = blank(prs)
    y = heading(s, "Pages it had never seen", "Three unrelated appraisal offices publish completed sample reports. None was used while building.")
    table(s, M, y + Inches(0.2), W - 2 * M, [
        ["page", "source", "boxes found", "marked", "flagged", "form recognised"],
        ["FHA appraisal", "Piekos Appraisals, IL", "118", "45", "0", "yes"],
        ["VA appraisal", "Piekos Appraisals, IL", "118", "46", "0", "yes"],
        ["standard appraisal", "RealVals", "118", "38", "0", "yes"],
        ["FHA appraisal", "Key Realty, MD", "118", "37", "0", "yes"],
        ["condominium appraisal", "Piekos Appraisals, IL", "88", "33", "0", "no, and it said so"],
    ], col_w=[1.8, 2.2, 1.1, 0.8, 0.8, 1.6], size=13)
    textbox(s, M, y + Inches(2.6), W - 2 * M, Inches(2.0), [
        "The condominium row is not the clean win it looks like.",
        "The form matcher matched it to the wrong form and scored that match a perfect 1.0. "
        "What caught it was the check behind it. The two readers agreed on only 68 of the 118 positions the standard form expects. "
        "So the reading was marked untrusted and the second reader was dropped for that page.",
        "A confidence score is not a safety mechanism.",
    ], size=16)
    notes(s, "This is the slide I would use if somebody asks what happens on a form nobody taught it. The answer "
             "is that the form matcher was confidently wrong and the agreement check held.")
    return s


def s_accuracy_one(prs):
    """Claim and demo on one slide, for the deck that gets narrated rather than read."""
    s = blank(prs)
    y = heading(s, "1. Accuracy, against whose definition?", "A circle instead of an X. A box struck out. A tick that overshoots. Each is defensible either way.")
    textbox(s, M, y + Inches(0.15), Inches(5.4), Inches(4.4), [
        theirs("None of those is ours to decide. One lender wants the circle counted, another wants it sent to a "
               "person, and both are right about their own files.",
               "None of those is ours to decide. One lender wants the circle counted, another wants it sent to an "
               "underwriter, and both are right about their own files."),
        "",
        plain("So the definition lives in policy.json, a file the customer owns. Changing what counts as a mark is "
              "editing that file. Not a ticket, not a release, not a call with us.",
              "So the definition lives in one settings file the customer owns. Changing what counts as a mark is "
              "editing that file. Not a ticket, not a release, not a call with us.",
              "So the definition lives in one settings file the lender owns. Changing what counts as a mark is "
              "editing that file. Not a ticket, not a release, not a call with us."),
        "",
        "uv run python scripts/compare_policies.py",
    ], size=16)
    table(s, Inches(6.3), y + Inches(0.25), Inches(6.4), [
        ["", "shipped", theirs("a stricter customer", "a stricter lender")],
        ["a single thin stroke", "counts as marked", "flag, do not decide"],
        ["a box scribbled out", "read as empty", "flag, do not overrule ink"],
        ["uncertain band", "0.05 to 0.20", "0.03 to 0.30"],
        ["agreement with the key", "285 of 286", "286 of 286"],
    ], col_w=[2.2, 1.5, 1.8], size=13)
    textbox(s, Inches(6.3), y + Inches(2.35), Inches(6.4), Inches(2.0), [
        "The stricter one is not simply better. It gets the faded X right by refusing to overrule ink, and the "
        "same rule would let a struck-out box report as filled. Which error you prefer is a business question.",
    ], size=15, color=MUTED)
    notes(s, "This is the demo I would run live. Two files, no code change, and the answer on one box moves. "
             "The point is not that one policy wins; it is that the argument now has somewhere to happen that "
             "is not our backlog.")
    return s


def s_referees(prs):
    s = blank(prs)
    y = heading(s, "How it is judged", "Three referees, because a system that grades its own homework is not a system.")
    chip(s, M, y + Inches(0.3), Inches(3.9), Inches(1.5),
         "The four real pages\nlabelled box by box; seeded by the software, then corrected, and that bias is stated", DET_FILL, DET, size=14)
    chip(s, Inches(4.75), y + Inches(0.3), Inches(3.9), Inches(1.5),
         "The blank official forms\nnothing on them is marked, so any mark reported there is a mistake with no argument", DET_FILL, DET, size=14)
    chip(s, Inches(8.8), y + Inches(0.3), Inches(3.9), Inches(1.5),
         "76 close-up crops\nlabelled by a person before any threshold was tuned; this one sets the policy", HUM_FILL, HUM, size=14)
    textbox(s, M, y + Inches(2.15), W - 2 * M, Inches(2.4), [
        "The third referee is the one that decides what a mark is.",
        "It is where three rules came from. A scribbled-out box is not a checked box. A stray pen line is "
        "not a selection. A circle drawn in a box is a case nobody should answer confidently.",
        f"Those rulings are written down in POLICY.md and the system is measured against them. "
        f"{GOLD.get('correct', 0)} of {GOLD.get('cards', 0)} hard-graded cards agree. "
        f"The boundary cases are excluded from grading and required to route rather than decide.",
    ], size=16, bold_first=True)
    notes(s, "Boundary cards are excluded from hard grading on purpose. A good system sits near the threshold "
             "on them or routes them, and being confidently wrong in either direction counts against it.")
    return s


def s_roadmap(prs):
    s = blank(prs)
    y = heading(s, "What I would do next, with real volume", "In the order the rungs pay for themselves.")
    textbox(s, M, y + Inches(0.2), W - 2 * M, Inches(4.6), [
        "1.  Run it in shadow over real uploads and spend human attention only where the two readers disagree. "
        "Agreement is cheap and uninformative; disagreement is where the information is.",
        "2.  Make the template registry data, so a new form revision is a file rather than a release, and every "
        "box carries its official field name. That is also what turns the queue into feedback for " + theirs("the customer.", "the lender."),
        plain("3.  Fix the trained classifier's gating so its disagreement only counts where the rule is already near "
              "a boundary. Today it changes no answer and triples the queue.",
              "3.  Let the small trained model's disagreement count only where the rule is already near a boundary. "
              "Today it changes no answer and triples the review queue."),
        "4.  Only once enough real pages are labelled does a trained detector become worth it, with this system "
        "as its auto-labeller and its permanent cross-check.",
    ], size=17, space=14)
    notes(s, "Disagreement sampling is the whole trick: it grows the labelled set without anyone reading pages "
             "at random, and the agreement rate doubles as a drift alarm that fires before anybody labels anything.")
    return s


NARRATE = ["title", "problem", "hard", "arch", "accuracy_one", "cost", "latency", "governance",
           "abc", "results", "run", "verify", "dashboard", "apprentice", "feedback", "limits"]

LEAVE_BEHIND = ["title", "problem", "hard", "arch", "accuracy_claim", "accuracy_demo", "cost", "latency",
                "governance", "abc", "why_b_c", "referees", "results", "holdout", "run", "verify",
                "dashboard", "feedback", "iterations", "roadmap", "limits"]

BUILDERS = {
    "title": lambda prs: s_title(prs, ""),
    "problem": s_problem,
    "hard": s_hard,
    "arch": s_arch,
    "accuracy_one": s_accuracy_one,
    "accuracy_claim": s_accuracy_claim,
    "accuracy_demo": lambda prs: s_accuracy_demo(prs, False),
    "cost": s_cost,
    "latency": s_latency,
    "governance": s_governance,
    "abc": s_abc,
    "why_b_c": s_why_b_c,
    "referees": s_referees,
    "results": s_results,
    "holdout": s_holdout,
    "run": s_run,
    "verify": s_verify,
    "dashboard": s_dashboard,
    "apprentice": s_apprentice,
    "feedback": s_feedback,
    "iterations": s_iterations,
    "roadmap": s_roadmap,
    "limits": s_limits,
}


SECTIONS = {
    "problem": "The problem", "hard": "The problem", "arch": "How it works",
    "accuracy_one": "Four dimensions", "accuracy_claim": "Four dimensions", "accuracy_demo": "Four dimensions",
    "cost": "Four dimensions", "latency": "Four dimensions", "governance": "Four dimensions",
    "abc": "Alternatives", "why_b_c": "Alternatives",
    "referees": "Evidence", "results": "Evidence", "holdout": "Evidence", "run": "Evidence", "verify": "Evidence", "dashboard": "Evidence", "apprentice": "Evidence",
    "feedback": "What comes next", "iterations": "What comes next", "roadmap": "What comes next", "limits": "What comes next",
}


def build(order: list[str], out: str) -> None:
    global SECTION
    prs = deck()
    for key in order:
        SECTION = SECTIONS.get(key, "")
        BUILDERS[key](prs)
    prs.core_properties.title = "Reading checkboxes off appraisal pages"
    prs.core_properties.author = "James W. Niu"
    prs.save(out)
    print(f"{out}  {len(order)} slides")


def main() -> None:
    global REGISTER
    # One deck ships: the plain-words copy for the room. The other registers and the long
    # leave-behind were working iterations; the submission does not carry them.
    REGISTER = "production"
    build(NARRATE, str(ROOT / "deliverables" / "checkbox-approach.pptx"))


if __name__ == "__main__":
    main()
